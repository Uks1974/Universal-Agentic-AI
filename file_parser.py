# file_parser.py

import pandas as pd
import pdfplumber
from docx import Document
from pptx import Presentation
from PIL import Image


def parse_file(file_path, file_type):
    """
    Parses uploaded files into text or structured data.
    OCR is intentionally disabled for Streamlit Cloud compatibility.
    """

    # -------- EXCEL --------
    if file_type == "excel":
        return pd.read_excel(file_path)

    # -------- PDF --------
    if file_type == "pdf":
        text = ""
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text += page.extract_text() or ""
        return text

    # -------- WORD --------
    if file_type == "word":
        doc = Document(file_path)
        return "\n".join([p.text for p in doc.paragraphs])

    # -------- POWERPOINT --------
    if file_type == "ppt":
        prs = Presentation(file_path)
        slides_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slides_text.append(shape.text)
        return "\n".join(slides_text)

    # -------- IMAGE (NO OCR) --------
    if file_type == "image":
        # OCR intentionally disabled on Streamlit Cloud
        return "Image uploaded. Text extraction (OCR) is disabled in this environment."

    return "Unsupported file type"
