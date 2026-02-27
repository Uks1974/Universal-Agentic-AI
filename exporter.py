# exporter.py
from pptx.util import Inches
from docx import Document
from pptx import Presentation
from reportlab.platypus import SimpleDocTemplate, Paragraph
from reportlab.lib.styles import getSampleStyleSheet
import pandas as pd
import matplotlib.pyplot as plt
import os

OUTPUT_DIR = "generated_files"
os.makedirs(OUTPUT_DIR, exist_ok=True)


def export_word(summary):
    doc = Document()
    doc.add_heading(summary["title"])
    for b in summary["bullets"]:
        doc.add_paragraph(b)
    path = os.path.join(OUTPUT_DIR, "Executive_Summary.docx")
    doc.save(path)
    return path


def export_excel(summary):
    df = pd.DataFrame({"Summary": summary["bullets"]})
    path = os.path.join(OUTPUT_DIR, "Executive_Summary.xlsx")
    df.to_excel(path, index=False)
    return path


def export_ppt(summary, chart_path=None):
    prs = Presentation()

    # Slide 1 – Executive Summary
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = summary["title"]
    slide.placeholders[1].text = "\n".join(summary["bullets"][:5])

    # Slide 2 – Recommendations & Actions
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Recommendations & Future Actions"
    slide2.placeholders[1].text = "\n".join(summary["bullets"][5:])

    # Slide 3 – Data Visual (if available)
    if chart_path and os.path.exists(chart_path):
        slide3 = prs.slides.add_slide(prs.slide_layouts[5])
        slide3.shapes.add_picture(chart_path, Inches(1), Inches(1), width=Inches(6))

    path = os.path.join(OUTPUT_DIR, "Executive_Summary.pptx")
    prs.save(path)
    return path


def export_pdf(summary):
    styles = getSampleStyleSheet()
    path = os.path.join(OUTPUT_DIR, "Executive_Summary.pdf")
    doc = SimpleDocTemplate(path)
    content = [Paragraph(summary["title"], styles["Title"])]
    for b in summary["bullets"]:
        content.append(Paragraph(b, styles["Normal"]))
    doc.build(content)
    return path


def export_image(summary):
    plt.figure(figsize=(8, 5))
    plt.text(0.01, 0.9, summary["title"], fontsize=16, weight="bold")
    for i, b in enumerate(summary["bullets"]):
        plt.text(0.01, 0.8 - i * 0.1, f"- {b}", fontsize=12)
    plt.axis("off")
    path = os.path.join(OUTPUT_DIR, "Executive_Summary.png")
    plt.savefig(path, bbox_inches="tight")
    plt.close()
    return path